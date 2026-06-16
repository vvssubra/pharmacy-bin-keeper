"""Build pharmacy-bin-keeper pitch deck — revised framing (drug monitoring & workflow)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── colour palette ──────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1B, 0x4F, 0x8A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x27, 0xAE, 0x60)
AMBER  = RGBColor(0xF3, 0x9C, 0x12)
RED    = RGBColor(0xE7, 0x4C, 0x3C)
LIGHT  = RGBColor(0xF0, 0xF4, 0xF8)
GREY   = RGBColor(0x4A, 0x55, 0x68)
DARK   = RGBColor(0x1A, 0x20, 0x2C)
BLUE2  = RGBColor(0x2D, 0x86, 0xC9)
PURPLE = RGBColor(0x8E, 0x44, 0xAD)
NAVY2  = RGBColor(0x12, 0x35, 0x62)
MUTED  = RGBColor(0xA8, 0xC4, 0xE8)
SBMUT  = RGBColor(0x90, 0xA8, 0xC0)
SMUT   = RGBColor(0xB0, 0xC8, 0xE8)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ── helpers ─────────────────────────────────────────────────────────────────
def rect(slide, x, y, w, h, fill):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.line.fill.background()
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    return sh


def tb(slide, text, x, y, w, h,
       size=20, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
       italic=False, wrap=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box


def bullets(slide, items, x, y, w, h, size=13, color=GREY, heading=None, hcolor=NAVY):
    box = slide.shapes.add_textbox(x, y, w, h)
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    if heading:
        p = tf.paragraphs[0]
        first = False
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = heading
        r.font.size = Pt(size + 2)
        r.font.bold = True
        r.font.color.rgb = hcolor
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        r = p.add_run()
        r.text = item
        r.font.size = Pt(size)
        r.font.color.rgb = color


def header(slide, title, subtitle=None):
    rect(slide, 0, 0, W, H, LIGHT)
    rect(slide, 0, 0, W, Inches(1.4), NAVY)
    tb(slide, title,    Inches(0.55), Inches(0.22), Inches(11), Inches(0.85),
       size=32, bold=True, color=WHITE)
    if subtitle:
        tb(slide, subtitle, Inches(0.55), Inches(0.88), Inches(11), Inches(0.48),
           size=15, color=SMUT)


def card(slide, cx, cy, cw, ch, fill=WHITE, bar_color=None):
    rect(slide, cx, cy, cw, ch, fill)
    if bar_color:
        rect(slide, cx, cy, cw, Inches(0.07), bar_color)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, NAVY)
rect(s, 0, 0, Inches(0.25), H, ACCENT)

tb(s, "Pharmacy Drug Monitoring System",
   Inches(0.65), Inches(1.6), Inches(10), Inches(1.3),
   size=44, bold=True, color=WHITE)
tb(s, "Digitising Quota Drug Requests & Antibiotic Forms\nfor Klinik Kesihatan",
   Inches(0.65), Inches(3.05), Inches(10), Inches(1.0),
   size=22, color=MUTED, italic=True, wrap=True)
tb(s, "Kawalan Ubat KK Kempas",
   Inches(0.65), Inches(4.2), Inches(7), Inches(0.5),
   size=16, color=ACCENT)
rect(s, 0, Inches(6.9), W, Inches(0.6), NAVY2)
tb(s, "Confidential · 2026", Inches(0.65), Inches(6.9), Inches(6), Inches(0.6),
   size=12, color=SBMUT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "The Problem",
       "How quota drug monitoring and antibiotic approval are handled today — and why it's breaking down")
rect(s, 0, 0, Inches(0.25), H, RED)

# Two main problem areas side by side
# LEFT — Quota drugs / Google Sheets
CX1, CX2 = Inches(0.55), Inches(6.95)
CW = Inches(6.0)

rect(s, CX1, Inches(1.55), CW, Inches(5.7), WHITE)
rect(s, CX1, Inches(1.55), CW, Inches(0.07), RED)
tb(s, "📊  Quota Drug Monitoring", CX1 + Inches(0.2), Inches(1.65),
   CW - Inches(0.3), Inches(0.6), size=17, bold=True, color=RED)
tb(s, "Currently tracked in Google Sheets",
   CX1 + Inches(0.2), Inches(2.3), CW - Inches(0.3), Inches(0.4),
   size=13, bold=True, color=DARK)
bullets(s, [
    "❌  Dozens of tabs — hard to find the right drug at the right time",
    "❌  Staff fill data into the wrong column or wrong tab",
    "❌  No real-time visibility — sheets go stale between updates",
    "❌  No workflow — approvals happen via WhatsApp or verbal",
    "❌  No audit trail — impossible to trace who approved what",
    "❌  MO, FMS, and pharmacist work from different versions",
],
CX1 + Inches(0.2), Inches(2.75), CW - Inches(0.3), Inches(4.2),
size=13, color=GREY)

# RIGHT — Antibiotic forms / physical paper
rect(s, CX2, Inches(1.55), CW, Inches(5.7), WHITE)
rect(s, CX2, Inches(1.55), CW, Inches(0.07), AMBER)
tb(s, "💉  Antibiotic Approval Forms", CX2 + Inches(0.2), Inches(1.65),
   CW - Inches(0.3), Inches(0.6), size=17, bold=True, color=AMBER)
tb(s, "Currently done on physical A4 photocopied sheets",
   CX2 + Inches(0.2), Inches(2.3), CW - Inches(0.3), Inches(0.4),
   size=13, bold=True, color=DARK)
bullets(s, [
    "❌  Printing and photocopying cost for every form",
    "❌  Paper forms lost, damaged, or illegible",
    "❌  Physical routing between MO → FMS → pharmacy is slow",
    "❌  No status visibility — MO doesn't know if form was approved",
    "❌  Storage & filing burden for compliance records",
    "❌  Cannot search historical forms quickly",
],
CX2 + Inches(0.2), Inches(2.75), CW - Inches(0.3), Inches(4.2),
size=13, color=GREY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — OUR SOLUTION
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Our Solution",
       "One system to replace Google Sheets and paper forms — with full workflow, visibility, and audit trail")
rect(s, 0, 0, Inches(0.25), H, ACCENT)

cards = [
    ("✅  Replace Google Sheets",  NAVY,
     "A structured digital register for every quota drug — searchable, real-time, and impossible to fill in the wrong place."),
    ("✅  Replace Paper Forms",    ACCENT,
     "Antibiotic Clinical Pathway forms (NAG 2024) submitted digitally. Zero printing cost, zero lost forms, instant routing."),
    ("✅  Structured Workflow",    BLUE2,
     "MO submits → FMS reviews & approves → Pharmacy acknowledges. Every step is tracked with a timestamp and a name."),
    ("✅  Real-Time Visibility",   AMBER,
     "All parties — MO, FMS, pharmacist — see the live status of every request. No chasing via WhatsApp."),
    ("✅  Full Audit Trail",       PURPLE,
     "Every approval, rejection, and acknowledgement is logged with the user's name, role, and timestamp. Compliance-ready."),
    ("✅  Lean & Paperless",       RGBColor(0x16, 0x7A, 0x6A),
     "Eliminate photocopy costs, A4 waste, filing cabinets, and manual data entry errors. A truly lean clinic operation."),
]

CW = Inches(3.9)
CH = Inches(2.3)
for i, (title, color, body) in enumerate(cards):
    col = i % 3
    row = i // 3
    cx = Inches(0.4) + col * Inches(4.28)
    cy = Inches(1.65) + row * Inches(2.55)
    card(s, cx, cy, CW, CH, WHITE, color)
    tb(s, title, cx + Inches(0.2), cy + Inches(0.18), CW - Inches(0.3), Inches(0.55),
       size=14, bold=True, color=color)
    tb(s, body,  cx + Inches(0.2), cy + Inches(0.75), CW - Inches(0.3), Inches(1.35),
       size=12, color=GREY, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — WORKFLOW A: QUOTA DRUG REQUEST
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Workflow 1 — Quota Drug Request",
       "MO issues a new quota drug to a patient: end-to-end digital, fully tracked")
rect(s, 0, 0, Inches(0.25), H, NAVY)
rect(s, 0, 0, W, H, WHITE)
rect(s, 0, 0, W, Inches(1.4), NAVY)
rect(s, 0, 0, Inches(0.25), H, AMBER)
tb(s, "Workflow 1 — Quota Drug Request", Inches(0.55), Inches(0.22),
   Inches(11), Inches(0.85), size=32, bold=True, color=WHITE)
tb(s, "MO issues a new quota drug to a patient: end-to-end digital, fully tracked",
   Inches(0.55), Inches(0.88), Inches(11), Inches(0.48), size=15, color=SMUT)

steps = [
    ("1", "MO Submits\nRequest",         NAVY,
     "Selects drug, enters\npatient IC & quantity.\nSubmits digitally."),
    ("2", "FMS Receives\n& Reviews",      AMBER,
     "Sees request on FMS\ndashboard. Reviews\npatient & drug details."),
    ("3", "FMS Approves\nor Rejects",     RED,
     "Approves with notes\nor rejects with reason.\nMO notified instantly."),
    ("4", "Pharmacy\nAcknowledges",       ACCENT,
     "Pharmacist sees approved\nrequest. Acknowledges\ndispensing to patient."),
    ("5", "Record\nSaved",                BLUE2,
     "Patient registry updated.\nFull audit log.\nQuota counter decrements."),
]

BW = Inches(2.1)
BH = Inches(4.2)
GAP = Inches(0.38)
SX = Inches(0.4)
SY = Inches(1.6)

for i, (num, title, color, body) in enumerate(steps):
    cx = SX + i * (BW + GAP)
    rect(s, cx, SY, BW, BH, LIGHT)
    rect(s, cx, SY, BW, Inches(0.07), color)
    # circle
    circ = s.shapes.add_shape(9, cx + Inches(0.65), SY + Inches(0.15), Inches(0.8), Inches(0.8))
    circ.fill.solid(); circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    tf = circ.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num
    r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = WHITE

    tb(s, title, cx + Inches(0.1), SY + Inches(1.05), BW - Inches(0.15), Inches(0.8),
       size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER, wrap=True)
    tb(s, body,  cx + Inches(0.1), SY + Inches(1.95), BW - Inches(0.15), Inches(2.0),
       size=12, color=GREY, align=PP_ALIGN.CENTER, wrap=True)

    if i < len(steps) - 1:
        ax = cx + BW + Inches(0.0)
        tb(s, "→", ax, SY + Inches(1.7), GAP, Inches(0.5),
           size=22, bold=True, color=GREY, align=PP_ALIGN.CENTER)

# bottom note
tb(s, "✓  Replaces: Google Sheets quota tracking with many tabs",
   Inches(0.5), Inches(6.8), Inches(12), Inches(0.45),
   size=13, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — WORKFLOW B: ANTIBIOTIC FORM
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, WHITE)
rect(s, 0, 0, W, Inches(1.4), NAVY)
rect(s, 0, 0, Inches(0.25), H, ACCENT)
tb(s, "Workflow 2 — Antibiotic Approval Form", Inches(0.55), Inches(0.22),
   Inches(11), Inches(0.85), size=32, bold=True, color=WHITE)
tb(s, "NAG 2024 Clinical Pathway form — digitised from A4 paper to the system",
   Inches(0.55), Inches(0.88), Inches(11), Inches(0.48), size=15, color=SMUT)

steps2 = [
    ("1", "MO Fills\nForm",              NAVY,
     "Enters patient, diagnosis,\nregimen & NAG 2024\nchecklist digitally."),
    ("2", "FMS Receives\n& Reviews",      AMBER,
     "Form appears on FMS\ndashboard. Reviews\nclinical details."),
    ("3", "FMS Approves\nor Rejects",     RED,
     "Approves or rejects\nwith specialist notes.\nDecision logged."),
    ("4", "Pharmacy\nAcknowledges",       ACCENT,
     "Pharmacist sees\napproved form and\nacknowledges dispensing."),
    ("5", "Record\nArchived",             BLUE2,
     "Digital record stored\nforever — searchable,\nno filing cabinets."),
]

for i, (num, title, color, body) in enumerate(steps2):
    cx = SX + i * (BW + GAP)
    rect(s, cx, SY, BW, BH, LIGHT)
    rect(s, cx, SY, BW, Inches(0.07), color)
    circ = s.shapes.add_shape(9, cx + Inches(0.65), SY + Inches(0.15), Inches(0.8), Inches(0.8))
    circ.fill.solid(); circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    tf = circ.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num
    r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = WHITE
    tb(s, title, cx + Inches(0.1), SY + Inches(1.05), BW - Inches(0.15), Inches(0.8),
       size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER, wrap=True)
    tb(s, body,  cx + Inches(0.1), SY + Inches(1.95), BW - Inches(0.15), Inches(2.0),
       size=12, color=GREY, align=PP_ALIGN.CENTER, wrap=True)
    if i < len(steps2) - 1:
        ax = cx + BW
        tb(s, "→", ax, SY + Inches(1.7), GAP, Inches(0.5),
           size=22, bold=True, color=GREY, align=PP_ALIGN.CENTER)

tb(s, "✓  Replaces: Physical A4 photocopy sheets — zero printing cost, zero lost forms",
   Inches(0.5), Inches(6.8), Inches(12), Inches(0.45),
   size=13, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — BEFORE vs AFTER
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, LIGHT)
rect(s, 0, 0, W, Inches(1.4), NAVY)
rect(s, 0, 0, Inches(0.25), H, BLUE2)
tb(s, "Before vs. After", Inches(0.55), Inches(0.22), Inches(11), Inches(0.85),
   size=32, bold=True, color=WHITE)
tb(s, "What changes for each team when they move to the system",
   Inches(0.55), Inches(0.88), Inches(11), Inches(0.48), size=15, color=SMUT)

# column headers
rect(s, Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.45), RED)
tb(s, "BEFORE  (Google Sheets + Paper)", Inches(0.6), Inches(1.5),
   Inches(5.6), Inches(0.45), size=14, bold=True, color=WHITE)
rect(s, Inches(6.9), Inches(1.5), Inches(6.0), Inches(0.45), ACCENT)
tb(s, "AFTER  (This System)", Inches(7.0), Inches(1.5),
   Inches(5.8), Inches(0.45), size=14, bold=True, color=WHITE)

rows = [
    ("Quota drug tracking",
     "Multiple Google Sheet tabs — easy to miss or enter in wrong tab",
     "Structured form per request — one place, always right"),
    ("Antibiotic approval",
     "Physical A4 photocopy form passed hand-to-hand",
     "Digital form submitted instantly, routed automatically"),
    ("FMS review",
     "Receive forms physically or via WhatsApp photo",
     "Dedicated FMS dashboard — all pending items in one view"),
    ("Pharmacy step",
     "Informed verbally or by paper after FMS approves",
     "Auto-notified with approved queue — one-click acknowledge"),
    ("Record keeping",
     "Files in folders; impossible to search quickly",
     "Searchable digital log — find any record in seconds"),
    ("Cost",
     "Photocopy & A4 paper cost for every antibiotic form",
     "Zero printing cost — fully paperless"),
]

for i, (aspect, before, after) in enumerate(rows):
    cy = Inches(2.05) + i * Inches(0.82)
    bg = WHITE if i % 2 == 0 else RGBColor(0xF7, 0xF9, 0xFC)
    rect(s, Inches(0.5), cy, Inches(12.4), Inches(0.78), bg)
    tb(s, aspect, Inches(0.6), cy + Inches(0.12), Inches(1.6), Inches(0.55),
       size=12, bold=True, color=NAVY)
    tb(s, "✗  " + before, Inches(2.3), cy + Inches(0.12), Inches(4.1), Inches(0.55),
       size=12, color=RED, wrap=True)
    tb(s, "✓  " + after,  Inches(6.6), cy + Inches(0.12), Inches(6.1), Inches(0.55),
       size=12, color=RGBColor(0x1A, 0x7A, 0x40), wrap=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — ROLES & WHAT EACH USER SEES
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, WHITE)
rect(s, 0, 0, W, Inches(1.4), NAVY)
rect(s, 0, 0, Inches(0.25), H, AMBER)
tb(s, "What Each Role Sees & Does", Inches(0.55), Inches(0.22), Inches(11), Inches(0.85),
   size=32, bold=True, color=WHITE)
tb(s, "Purpose-built view for each team — no clutter, no confusion",
   Inches(0.55), Inches(0.88), Inches(11), Inches(0.48), size=15, color=SMUT)

roles_data = [
    ("👨‍⚕️  Medical Officer\n(MO)", NAVY, [
        "Submit quota drug requests for patients",
        "Fill and submit antibiotic approval forms",
        "Track live status of submitted requests",
        "View pending count badge on dashboard",
    ]),
    ("🩺  FMS Specialist", AMBER, [
        "Dedicated inbox for all pending MO requests",
        "Review quota drug requests & approve / reject",
        "Review antibiotic forms & approve / reject",
        "Add specialist notes to every decision",
        "Monitor annual quota usage per drug",
    ]),
    ("💊  Pharmacist /\nAdmin", ACCENT, [
        "Receive approved requests automatically",
        "Acknowledge dispensing (quota drugs)",
        "Acknowledge antibiotic forms post-FMS",
        "Manage drug master list & stock receipts",
        "View reports and patient drug records",
    ]),
]

CW = Inches(3.9)
CH = Inches(5.6)
for i, (role, color, items) in enumerate(roles_data):
    cx = Inches(0.45) + i * Inches(4.28)
    card(s, cx, Inches(1.55), CW, CH, LIGHT, color)
    tb(s, role, cx + Inches(0.2), Inches(1.75), CW - Inches(0.3), Inches(0.75),
       size=16, bold=True, color=color, wrap=True)
    bullets(s, ["• " + it for it in items],
            cx + Inches(0.2), Inches(2.6), CW - Inches(0.3), CH - Inches(1.15),
            size=13, color=GREY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — KEY BENEFITS (LEAN)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, LIGHT)
rect(s, 0, 0, W, Inches(1.4), NAVY)
rect(s, 0, 0, Inches(0.25), H, ACCENT)
tb(s, "Key Benefits", Inches(0.55), Inches(0.22), Inches(11), Inches(0.85),
   size=32, bold=True, color=WHITE)
tb(s, "Measurable improvements from day one of going live",
   Inches(0.55), Inches(0.88), Inches(11), Inches(0.48), size=15, color=SMUT)

benefits = [
    ("💸", "Cost Reduction",         ACCENT,
     "Eliminate photocopy and A4 paper costs for every antibiotic form. Savings compound across the year."),
    ("⚡", "Speed",                   AMBER,
     "Requests reach FMS the moment MO submits. No walking, no waiting. FMS approves; pharmacy is notified instantly."),
    ("🎯", "Accuracy",                NAVY,
     "Structured forms mean no more wrong columns or wrong tabs. Every field is validated before submission."),
    ("🔍", "Searchability",           BLUE2,
     "Find any quota drug request or antibiotic form in seconds by patient IC, drug name, or date."),
    ("📜", "Compliance & Audit",      PURPLE,
     "Every action is attributed to a named user with a timestamp. Full audit trail for MOH inspections."),
    ("📊", "Quota Visibility",        RED,
     "Real-time annual quota counter per drug. FMS and pharmacist always know how much quota remains."),
]

CW = Inches(3.9)
CH = Inches(2.3)
for i, (icon, title, color, body) in enumerate(benefits):
    col = i % 3
    row = i // 3
    cx = Inches(0.4) + col * Inches(4.28)
    cy = Inches(1.65) + row * Inches(2.55)
    card(s, cx, cy, CW, CH, WHITE, color)
    tb(s, icon + "  " + title, cx + Inches(0.2), cy + Inches(0.18),
       CW - Inches(0.3), Inches(0.55), size=15, bold=True, color=color)
    tb(s, body, cx + Inches(0.2), cy + Inches(0.75), CW - Inches(0.3), Inches(1.35),
       size=12, color=GREY, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — TECH STACK
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, WHITE)
rect(s, 0, 0, W, Inches(1.4), NAVY)
rect(s, 0, 0, Inches(0.25), H, PURPLE)
tb(s, "Technology Stack", Inches(0.55), Inches(0.22), Inches(11), Inches(0.85),
   size=32, bold=True, color=WHITE)
tb(s, "Modern, cloud-native — zero servers to manage, accessible from any browser",
   Inches(0.55), Inches(0.88), Inches(11), Inches(0.48), size=15, color=SMUT)

layers = [
    ("Frontend",        NAVY,   ["React 18 + TypeScript", "Vite · React Router v6", "TanStack Query (real-time)", "Auto-refresh every 15 – 30 s"]),
    ("UI / Design",     BLUE2,  ["shadcn/ui (Radix)", "Tailwind CSS", "Recharts (quota graphs)", "Malay-language UI"]),
    ("Backend / Data",  ACCENT, ["Supabase (PostgreSQL)", "Row-Level Security (RLS)", "Supabase Auth (JWT)", "Cloud-hosted, auto-backup"]),
    ("Quality",         AMBER,  ["React Hook Form + Zod", "Vitest unit tests", "Playwright E2E tests", "TypeScript strict mode"]),
]

CW = Inches(2.9)
for i, (layer, color, items) in enumerate(layers):
    cx = Inches(0.5) + i * Inches(3.1)
    cy = Inches(1.6)
    ch = Inches(5.6)
    card(s, cx, cy, CW, ch, LIGHT)
    rect(s, cx, cy, CW, Inches(0.5), color)
    tb(s, layer, cx + Inches(0.15), cy + Inches(0.07), CW - Inches(0.3), Inches(0.4),
       size=15, bold=True, color=WHITE)
    bullets(s, ["▸  " + it for it in items],
            cx + Inches(0.2), cy + Inches(0.65), CW - Inches(0.3), ch - Inches(0.75),
            size=13, color=GREY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — ROADMAP
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, LIGHT)
rect(s, 0, 0, W, Inches(1.4), NAVY)
rect(s, 0, 0, Inches(0.25), H, ACCENT)
tb(s, "Product Roadmap", Inches(0.55), Inches(0.22), Inches(11), Inches(0.85),
   size=32, bold=True, color=WHITE)
tb(s, "From single KK to network-wide rollout",
   Inches(0.55), Inches(0.88), Inches(11), Inches(0.48), size=15, color=SMUT)

phases = [
    ("✅  Phase 1\n(Done)", ACCENT, [
        "Quota drug request workflow",
        "Antibiotic form (NAG 2024) workflow",
        "MO → FMS → Pharmacy approval chain",
        "Annual quota monitoring per drug",
        "Patient registry & drug history",
        "Full audit trail & reporting",
        "Role-based dashboards (MO / FMS / Pharmacist)",
    ]),
    ("🔄  Phase 2\n(In Progress)", AMBER, [
        "PDF export for antibiotic forms",
        "Push / in-app notifications on approval",
        "Advanced laporan (XLSX export)",
        "Drug interaction alerts",
        "Multi-facility support (facility switcher)",
    ]),
    ("🚀  Phase 3\n(Planned)", NAVY, [
        "MOH / JKN system integration",
        "AI-powered quota forecasting",
        "Mobile-friendly PWA wrapper",
        "Network-wide analytics across KK branches",
        "Batch antibiotic form review for FMS",
    ]),
]

CW = Inches(3.7)
for i, (phase, color, items) in enumerate(phases):
    cx = Inches(0.5) + i * Inches(4.1)
    cy = Inches(1.6)
    ch = Inches(5.6)
    card(s, cx, cy, CW, ch, WHITE, color)
    tb(s, phase, cx + Inches(0.2), cy + Inches(0.18), CW - Inches(0.3), Inches(0.75),
       size=16, bold=True, color=color, wrap=True)
    bullets(s, ["• " + it for it in items],
            cx + Inches(0.2), cy + Inches(1.0), CW - Inches(0.3), ch - Inches(1.1),
            size=12, color=GREY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — CLOSING
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, NAVY)
rect(s, 0, 0, Inches(0.25), H, ACCENT)

tb(s, "Pharmacy Drug Monitoring System",
   Inches(0.65), Inches(1.4), Inches(10.5), Inches(1.2),
   size=40, bold=True, color=WHITE)
tb(s, "Replacing Google Sheets and paper forms with a lean,\ndigital workflow — for every Klinik Kesihatan.",
   Inches(0.65), Inches(2.75), Inches(10), Inches(1.2),
   size=21, color=MUTED, italic=True, wrap=True)

# two impact statements
rect(s, Inches(0.65), Inches(4.2), Inches(5.5), Inches(0.95), RGBColor(0x17, 0x3A, 0x6E))
tb(s, "📊  Quota drugs → off Google Sheets, onto one system",
   Inches(0.85), Inches(4.35), Inches(5.2), Inches(0.65),
   size=15, color=ACCENT, bold=True, wrap=True)

rect(s, Inches(6.8), Inches(4.2), Inches(6.0), Inches(0.95), RGBColor(0x17, 0x3A, 0x6E))
tb(s, "💉  Antibiotic forms → off A4 paper, zero printing cost",
   Inches(7.0), Inches(4.35), Inches(5.7), Inches(0.65),
   size=15, color=AMBER, bold=True, wrap=True)

tb(s, "Contact  ·  admin@vvsdigitalsolutions.com",
   Inches(0.65), Inches(5.7), Inches(9), Inches(0.55),
   size=16, color=ACCENT)
rect(s, 0, Inches(6.9), W, Inches(0.6), NAVY2)
tb(s, "vvsdigitalsolutions.com  ·  Confidential 2026",
   Inches(0.65), Inches(6.9), Inches(10), Inches(0.6),
   size=12, color=SBMUT)


# ── Save ─────────────────────────────────────────────────────────────────────
out = "/home/user/pharmacy-bin-keeper/docs/pharmacy-bin-keeper-deck.pptx"
prs.save(out)
print(f"Saved → {out}")
